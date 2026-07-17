"""
Document generation router.
Generates DOCX (python-docx), PPTX (python-pptx), and PDF (DOCX → LibreOffice).
"""
import io
import os
import re
import tempfile
import subprocess
from typing import List, Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from pydantic import BaseModel

from app.core.deps import get_current_user
from app.database import get_db
from app.models.user import User

router = APIRouter(prefix="/api/v1/generate", tags=["generate"])


# ── Pydantic models ────────────────────────────────────────────────────────────

class DocxSection(BaseModel):
    heading: str
    content: str
    level: int = 1  # 1–4


class DocxRequest(BaseModel):
    title: str
    sections: List[DocxSection] = []
    author: Optional[str] = None
    font: str = "Calibri"
    font_size: int = 11


class PptxSlide(BaseModel):
    title: str
    bullets: List[str] = []
    notes: Optional[str] = None


class PptxRequest(BaseModel):
    title: str
    subtitle: Optional[str] = None
    slides: List[PptxSlide] = []
    theme: str = "default"  # default | dark | minimal


class PdfRequest(BaseModel):
    title: str
    content: str  # markdown-ish
    author: Optional[str] = None
    page_size: str = "A4"  # A4 | Letter | A3


# ── DOCX ───────────────────────────────────────────────────────────────────────

@router.post("/docx")
async def generate_docx(
    body: DocxRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """Generate a Word document with title, author, and dynamic sections."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    VALID_FONTS = {"Calibri", "Arial", "Times New Roman", "Georgia"}
    font_name = body.font if body.font in VALID_FONTS else "Calibri"

    doc = Document()

    # Set document properties
    core_props = doc.core_properties
    core_props.title = body.title
    if body.author:
        core_props.author = body.author

    # Title paragraph
    title_para = doc.add_heading(body.title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title_para.runs:
        run.font.name = font_name

    if body.author:
        author_para = doc.add_paragraph(f"Author: {body.author}")
        author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in author_para.runs:
            run.font.name = font_name
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.add_paragraph()  # spacer

    # Sections
    for section in body.sections:
        level = max(1, min(4, section.level))
        heading = doc.add_heading(section.heading, level=level)
        for run in heading.runs:
            run.font.name = font_name

        # Split content into paragraphs
        for paragraph_text in section.content.split("\n"):
            paragraph_text = paragraph_text.strip()
            if not paragraph_text:
                continue
            if paragraph_text.startswith("- ") or paragraph_text.startswith("* "):
                p = doc.add_paragraph(paragraph_text[2:], style="List Bullet")
            elif re.match(r"^\d+\.\s", paragraph_text):
                p = doc.add_paragraph(re.sub(r"^\d+\.\s", "", paragraph_text), style="List Number")
            else:
                p = doc.add_paragraph(paragraph_text)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(body.font_size)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    filename = _safe_filename(body.title) + ".docx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── PPTX ───────────────────────────────────────────────────────────────────────

@router.post("/pptx")
async def generate_pptx(
    body: PptxRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """Generate a PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Theme color definitions
    THEMES = {
        "default": {
            "bg": RGBColor(0xFF, 0xFF, 0xFF),
            "title_text": RGBColor(0x1A, 0x1A, 0x2E),
            "body_text": RGBColor(0x33, 0x33, 0x33),
            "accent": RGBColor(0x16, 0x4D, 0xAA),
        },
        "dark": {
            "bg": RGBColor(0x1A, 0x1A, 0x2E),
            "title_text": RGBColor(0xFF, 0xFF, 0xFF),
            "body_text": RGBColor(0xCC, 0xCC, 0xCC),
            "accent": RGBColor(0x4A, 0x9E, 0xFF),
        },
        "minimal": {
            "bg": RGBColor(0xF8, 0xF8, 0xF6),
            "title_text": RGBColor(0x11, 0x11, 0x11),
            "body_text": RGBColor(0x44, 0x44, 0x44),
            "accent": RGBColor(0x22, 0x22, 0x22),
        },
    }

    theme = THEMES.get(body.theme, THEMES["default"])

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def _set_bg(slide, color: RGBColor):
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(slide, text, left, top, width, height, font_size, bold=False, color=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    # Title slide
    blank_layout = prs.slide_layouts[6]  # blank layout
    title_slide = prs.slides.add_slide(blank_layout)
    _set_bg(title_slide, theme["bg"])

    # Title
    _add_text_box(
        title_slide, body.title,
        Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
        font_size=40, bold=True, color=theme["title_text"], align=PP_ALIGN.CENTER,
    )

    # Subtitle
    if body.subtitle:
        _add_text_box(
            title_slide, body.subtitle,
            Inches(1), Inches(4.2), Inches(11.33), Inches(0.8),
            font_size=24, bold=False, color=theme["body_text"], align=PP_ALIGN.CENTER,
        )

    # Content slides
    for slide_data in body.slides:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, theme["bg"])

        # Slide title
        _add_text_box(
            slide, slide_data.title,
            Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.0),
            font_size=28, bold=True, color=theme["accent"],
        )

        # Divider line (thin rectangle)
        from pptx.util import Inches, Pt, Emu
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.5), Inches(1.4), Inches(12.33), Emu(18000),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme["accent"]
        line.line.fill.background()

        # Bullets
        if slide_data.bullets:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(slide_data.bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.space_before = Pt(4)
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = theme["body_text"]

        # Speaker notes
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.notes

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = _safe_filename(body.title) + ".pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── PDF ────────────────────────────────────────────────────────────────────────

@router.post("/pdf")
async def generate_pdf(
    body: PdfRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """
    Generate a PDF from markdown-ish content.
    Strategy: build DOCX with python-docx, then convert with LibreOffice headless.
    """
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    # Page size mapping (width x height in inches)
    PAGE_SIZES = {
        "A4":     (8.27, 11.69),
        "Letter": (8.5,  11.0),
        "A3":     (11.69, 16.54),
    }
    page_w, page_h = PAGE_SIZES.get(body.page_size, PAGE_SIZES["A4"])

    doc = Document()

    # Set page size
    section = doc.sections[0]
    section.page_width  = Inches(page_w)
    section.page_height = Inches(page_h)
    section.left_margin   = Inches(1.0)
    section.right_margin  = Inches(1.0)
    section.top_margin    = Inches(1.0)
    section.bottom_margin = Inches(1.0)

    # Document properties
    doc.core_properties.title = body.title
    if body.author:
        doc.core_properties.author = body.author

    # Title
    title_para = doc.add_heading(body.title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if body.author:
        a = doc.add_paragraph(f"Author: {body.author}")
        a.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in a.runs:
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            run.font.size = Pt(11)

    doc.add_paragraph()

    # Parse markdown content line by line
    lines = body.content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]

        # ATX headings
        if line.startswith("#### "):
            doc.add_heading(line[5:].strip(), level=4)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        # Unordered list
        elif line.startswith("- ") or line.startswith("* "):
            p = doc.add_paragraph(line[2:].strip(), style="List Bullet")
            for run in p.runs:
                run.font.size = Pt(11)
        # Ordered list
        elif re.match(r"^\d+\.\s", line):
            text = re.sub(r"^\d+\.\s", "", line).strip()
            p = doc.add_paragraph(text, style="List Number")
            for run in p.runs:
                run.font.size = Pt(11)
        # Horizontal rule
        elif line.strip() in ("---", "***", "___"):
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), "AAAAAA")
            pBdr.append(bottom)
            pPr.append(pBdr)
        # Blockquote
        elif line.startswith("> "):
            p = doc.add_paragraph(line[2:].strip(), style="Quote")
            for run in p.runs:
                run.font.size = Pt(11)
        # Bold/italic inline (basic)
        elif line.strip() == "":
            doc.add_paragraph()
        else:
            p = doc.add_paragraph()
            _add_inline_paragraph(p, line.strip())

        i += 1

    # Save DOCX to temp file, convert to PDF with LibreOffice
    with tempfile.TemporaryDirectory() as tmp:
        docx_path = os.path.join(tmp, "document.docx")
        doc.save(docx_path)

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", tmp,
            docx_path,
        ]

        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120,
            )
        except subprocess.TimeoutExpired:
            raise HTTPException(status_code=500, detail="LibreOffice conversion timed out")

        if result.returncode != 0:
            raise HTTPException(
                status_code=500,
                detail=f"PDF conversion failed: {result.stderr[-500:] if result.stderr else 'unknown'}",
            )

        pdf_path = os.path.join(tmp, "document.pdf")
        if not os.path.exists(pdf_path):
            raise HTTPException(status_code=500, detail="PDF file was not created")

        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

    filename = _safe_filename(body.title) + ".pdf"
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Helpers ────────────────────────────────────────────────────────────────────

def _safe_filename(name: str) -> str:
    """Sanitize a string for use as a filename."""
    name = re.sub(r'[^\w\s\-]', '', name)
    name = re.sub(r'\s+', '_', name.strip())
    return name[:64] or "document"


def _add_inline_paragraph(paragraph, text: str):
    """Add a paragraph with basic **bold** and *italic* markdown inline support."""
    from docx.shared import Pt

    # Pattern to split on **bold** and *italic*
    pattern = re.compile(r'(\*\*[^*]+\*\*|\*[^*]+\*)')
    parts = pattern.split(text)

    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
            run.font.size = Pt(11)
        elif part.startswith("*") and part.endswith("*"):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
            run.font.size = Pt(11)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(11)
